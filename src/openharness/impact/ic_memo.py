"""IC (Investment Committee) memo generator.

Produces a structured IC memo from an `Assessment`, an `IC gate scorecard`,
and an optional `FundThesis`. Outputs in Markdown by default; if optional
deps are installed (`python-docx`, `python-pptx`), the same content is
emitted as Word / PowerPoint.

The Markdown is intentionally LP-friendly and follows a format that maps to
each common GP IC template (Sequoia / TPG Rise / generic), so a GP can paste
the output directly or pipe it through `pandoc` to their preferred format.
"""
from __future__ import annotations

from datetime import date
from html import escape
from pathlib import Path
from typing import Literal

from openharness.impact.deal_gate import DealScorecard
from openharness.impact.fund_thesis import FundThesis, weighted_5d_overall, weighted_sdg_overall
from openharness.impact.models import Assessment
from openharness.impact.report_templates.report_v2 import (
    GRADE_CLASS,
    render_footer,
    render_hero,
    render_kpi_strip,
    render_toc,
    sdg_swatch,
    wrap_document,
)


def _fmt_pct(x: float) -> str:
    return f"{x:.1f}%"


def _fmt_score(x: float) -> str:
    return f"{x:.1f}/5"


def render_ic_memo_markdown(
    assessment: Assessment,
    scorecard: DealScorecard,
    thesis: FundThesis | None = None,
    *,
    dd_coverage_pct: float | None = None,
    greenwashing_score: float | None = None,
    greenwashing_classification: str | None = None,
    deal_size_eur_m: float | None = None,
) -> str:
    """Build a Markdown IC memo string."""
    company = assessment.company
    fd = assessment.five_dimensions
    sdgs = sorted(assessment.sdg_alignments, key=lambda a: a.score, reverse=True)
    today = date.today().isoformat()

    out: list[str] = []
    out.append(f"# Investment Committee Memo — {company.name or 'Unnamed Company'}")
    out.append("")
    out.append(f"- **Date**: {today}")
    out.append(f"- **Fund**: {scorecard.fund}")
    out.append(f"- **Sector**: {company.sector or 'n/a'}")
    if deal_size_eur_m is not None:
        out.append(f"- **Proposed deal size**: EUR {deal_size_eur_m:.1f}m")
    out.append(f"- **IC gate result**: **{scorecard.overall_status.upper()}** — {scorecard.recommendation}")
    out.append("")

    out.append("## 1. Impact Thesis Fit")
    if thesis and not thesis.is_default:
        out.append(f"Fund thesis: *{thesis.name}* ({thesis.strategy}, vintage {thesis.vintage_year})")
        if fd:
            w5d = weighted_5d_overall(fd, thesis)
            out.append(f"- Fund-weighted 5D overall: **{w5d:.2f}/5** (unweighted: {_fmt_score(fd.overall_score)})")
        if sdgs:
            wsdg = weighted_sdg_overall(sdgs, thesis)
            out.append(f"- Fund-weighted SDG score: **{wsdg:.1f}/100**")
    else:
        out.append("*No fund thesis loaded — using equal-weight defaults. Set `data/fund_thesis.yaml` for fund-specific weighting.*")
    out.append("")

    out.append("## 2. 5-Dimension Impact Score")
    if fd:
        out.append("| Dimension | Score | Provenance | Notes |")
        out.append("|---|---|---|---|")
        for dim_name, ds in (
            ("What", fd.what), ("Who", fd.who), ("How Much", fd.how_much),
            ("Contribution", fd.contribution), ("Risk", fd.risk),
        ):
            out.append(f"| {dim_name} | {_fmt_score(ds.score)} | {ds.provenance} | {ds.notes} |")
        out.append(f"| **Overall** | **{_fmt_score(fd.overall_score)}** | {fd.overall_provenance} | Grade: {fd.overall_grade} |")
        if fd.recommendations:
            out.append("")
            out.append("**Recommendations:**")
            for r in fd.recommendations:
                out.append(f"- {r}")
    else:
        out.append("*No 5D assessment available.*")
    out.append("")

    out.append("## 3. SDG Alignment (top 5)")
    if sdgs:
        out.append("| SDG | Score | Targets matched | Provenance |")
        out.append("|---|---|---|---|")
        for a in sdgs[:5]:
            out.append(
                f"| SDG {a.goal} ({a.goal_name}) | {a.score:.0f}/100 "
                f"| {len(a.matched_targets)} | {a.provenance} ({a.scoring_basis}) |"
            )
    out.append("")

    out.append("## 4. Due Diligence Coverage")
    if dd_coverage_pct is not None:
        out.append(f"- Coverage: **{_fmt_pct(dd_coverage_pct)}**")
    else:
        out.append("- Coverage: *not run*")
    out.append("")

    out.append("## 5. Greenwashing & Authenticity")
    if greenwashing_score is not None:
        out.append(f"- Composite risk: **{greenwashing_score:.1f}/100**"
                   + (f" — {greenwashing_classification}" if greenwashing_classification else ""))
    else:
        out.append("- *Greenwashing screen not run.*")
    out.append("")

    out.append("## 6. IC Gate Detail")
    out.append("| Check | Status | Actual | Threshold | Note |")
    out.append("|---|---|---|---|---|")
    for c in scorecard.checks:
        out.append(f"| {c.name} | {c.status} | {c.actual} | {c.threshold} | {c.message} |")
    out.append("")

    if scorecard.blocking_failures:
        out.append("### Blocking failures")
        for m in scorecard.blocking_failures:
            out.append(f"- ❌ {m}")
        out.append("")
    if scorecard.warnings_list:
        out.append("### Warnings")
        for m in scorecard.warnings_list:
            out.append(f"- ⚠️ {m}")
        out.append("")

    out.append("## 7. Recommendation")
    out.append(f"> {scorecard.recommendation}")
    out.append("")
    out.append("---")
    out.append("*Generated by Impact Vision — review before submission.*")

    return "\n".join(out)


def render_ic_memo_docx(markdown: str, path: str | Path) -> Path:
    """Render an IC memo as a .docx file.

    Requires `python-docx`. If not installed, raises ImportError with a clear
    install hint. Heading detection is intentionally simple — production GPs
    will pipe the markdown through pandoc for richer formatting.
    """
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:
        raise ImportError(
            "Word export requires the optional dependency 'python-docx'. "
            "Install with: pip install python-docx"
        ) from exc

    doc = Document()
    for line in markdown.splitlines():
        if line.startswith("# "):
            doc.add_heading(line[2:], level=0)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=1)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=2)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.startswith("> "):
            doc.add_paragraph(line[2:], style="Intense Quote")
        elif line.startswith("|"):
            # leave tables as plain text — pandoc / a richer renderer handles
            # actual table layout. Avoid silent formatting bugs.
            doc.add_paragraph(line)
        else:
            doc.add_paragraph(line)

    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(p))
    return p


def render_ic_memo_pptx(
    assessment: Assessment,
    scorecard: DealScorecard,
    thesis: FundThesis | None = None,
    path: str | Path = "ic_memo.pptx",
) -> Path:
    """Render a 4-slide IC summary deck.

    Requires `python-pptx`.
    """
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches
    except ImportError as exc:
        raise ImportError(
            "PowerPoint export requires 'python-pptx'. Install with: pip install python-pptx"
        ) from exc

    prs = Presentation()

    # Slide 1: Title
    blank = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.title.text = f"IC Memo — {assessment.company.name or 'Unnamed'}"

    # Slide 2: 5D
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "5-Dimension Impact Score"
    if assessment.five_dimensions:
        fd = assessment.five_dimensions
        body = (f"Overall: {fd.overall_score:.1f}/5 ({fd.overall_grade})\n"
                f"What: {fd.what.score} | Who: {fd.who.score} | "
                f"How Much: {fd.how_much.score} | "
                f"Contribution: {fd.contribution.score} | Risk: {fd.risk.score}")
        s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2)).text_frame.text = body

    # Slide 3: SDG
    s3 = prs.slides.add_slide(blank)
    s3.shapes.title.text = "SDG Alignment"
    sdg_lines = "\n".join(
        f"SDG {a.goal} ({a.goal_name}): {a.score:.0f}/100"
        for a in sorted(assessment.sdg_alignments, key=lambda a: a.score, reverse=True)[:5]
    )
    s3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4)).text_frame.text = sdg_lines

    # Slide 4: IC gate
    s4 = prs.slides.add_slide(blank)
    s4.shapes.title.text = f"IC Gate: {scorecard.overall_status.upper()}"
    body = "\n".join(f"- {c.name}: {c.status}" for c in scorecard.checks)
    body += "\n\nRecommendation:\n" + scorecard.recommendation
    s4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4)).text_frame.text = body

    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return p


def _status_pill(status: str) -> str:
    cls = {"pass": "pass", "warn": "warn", "fail": "fail"}.get(status.lower(), "na")
    return f'<span class="pill {cls}">{escape(status.upper())}</span>'


def _gw_kind(score: float | None) -> tuple[str, str]:
    """Return (tile-kind, badge-label) for a greenwashing composite score."""
    if score is None:
        return "neutral", "not run"
    if score >= 70:
        return "fail", "high risk"
    if score >= 40:
        return "warn", "elevated"
    return "pass", "low risk"


def render_ic_memo_html(
    assessment: Assessment,
    scorecard: DealScorecard,
    thesis: FundThesis | None = None,
    *,
    dd_coverage_pct: float | None = None,
    greenwashing_score: float | None = None,
    greenwashing_classification: str | None = None,
    deal_size_eur_m: float | None = None,
) -> str:
    """Render the IC memo as a self-contained, print-ready HTML document.

    Uses the shared v2 chrome (hero + KPI strip + sticky TOC + callouts)
    so the output can be shared with LPs or archived in the signed feed
    without depending on Markdown renderers.
    """
    company = assessment.company
    fd = assessment.five_dimensions
    sdgs = sorted(assessment.sdg_alignments, key=lambda a: a.score, reverse=True)
    today = date.today().isoformat()
    top = sdgs[0] if sdgs else None

    # --- Hero -------------------------------------------------------
    meta = [("Date", today), ("Fund", scorecard.fund)]
    if company.sector:
        meta.append(("Sector", company.sector))
    if getattr(company, "geography", ""):
        meta.append(("Geography", company.geography))
    if deal_size_eur_m is not None:
        meta.append(("Deal size", f"EUR {deal_size_eur_m:.1f}m"))
    meta.append(("IC gate", scorecard.overall_status.upper()))

    hero = render_hero(
        eyebrow="Investment Committee Memo",
        title=escape(company.name or "Unnamed Company"),
        subtitle=escape(scorecard.recommendation or ""),
        meta=meta,
        tags=(company.impact_themes or [])[:6],
    )

    # --- KPI strip --------------------------------------------------
    kpis: list[dict] = []
    gate_kind = {"pass": "pass", "warn": "warn", "fail": "fail"}.get(
        scorecard.overall_status, "neutral",
    )
    kpis.append({
        "label": "IC Gate",
        "value": scorecard.overall_status.upper(),
        "sub": f"{len(scorecard.checks)} checks · "
               f"{len(scorecard.blocking_failures)} fail / {len(scorecard.warnings_list)} warn",
        "kind": gate_kind,
    })
    if fd:
        kpis.append({
            "label": "5D Overall",
            "value": f'<span class="{GRADE_CLASS.get(fd.overall_grade, "")}">'
                     f'{fd.overall_score:.1f}<span style="font-size:0.55em;color:var(--text-muted)">/5</span></span>',
            "sub": f"Grade {fd.overall_grade} · {fd.overall_provenance}",
            "kind": "pass" if fd.overall_score >= 3.5 else "warn" if fd.overall_score >= 2.5 else "neutral",
        })
    if top:
        kpis.append({
            "label": "Top SDG",
            "value": f"SDG {top.goal}",
            "sub": f"{top.goal_name} · {top.score:.0f}/100",
            "kind": "pass" if top.score >= 60 else "warn" if top.score >= 30 else "neutral",
        })
    if dd_coverage_pct is not None:
        kpis.append({
            "label": "DD Coverage",
            "value": f"{dd_coverage_pct:.0f}%",
            "sub": "Checklist addressed",
            "kind": "pass" if dd_coverage_pct >= 70 else "warn" if dd_coverage_pct >= 40 else "fail",
        })
    gw_kind, gw_badge = _gw_kind(greenwashing_score)
    kpis.append({
        "label": "Greenwashing risk",
        "value": f"{greenwashing_score:.0f}/100" if greenwashing_score is not None else "—",
        "sub": greenwashing_classification or "",
        "badge": gw_badge, "badge_kind": gw_kind, "kind": gw_kind,
    })

    # --- Section: Thesis fit ---------------------------------------
    thesis_html: list[str] = ['<section class="card" id="thesis-fit">'
                              '<h2 class="section-title">Impact thesis fit</h2>']
    if thesis and not thesis.is_default:
        thesis_html.append(
            f'<p class="section-lede">Fund thesis: <b>{escape(thesis.name)}</b> '
            f'({escape(thesis.strategy or "")}, vintage {thesis.vintage_year}).</p>',
        )
        thesis_html.append('<div class="cards-row">')
        if fd:
            w5d = weighted_5d_overall(fd, thesis)
            thesis_html.append(
                f'<div class="score-card"><div class="value">{w5d:.2f}<span '
                f'style="font-size:0.45em;color:var(--text-muted)">/5</span></div>'
                f'<div class="label">Fund-weighted 5D</div></div>')
        if sdgs:
            wsdg = weighted_sdg_overall(sdgs, thesis)
            thesis_html.append(
                f'<div class="score-card"><div class="value">{wsdg:.0f}<span '
                f'style="font-size:0.45em;color:var(--text-muted)">/100</span></div>'
                f'<div class="label">Fund-weighted SDG</div></div>')
        thesis_html.append('</div>')
    else:
        thesis_html.append(
            '<div class="callout warn">No fund thesis loaded — using equal-weight defaults. '
            'Set <code>data/fund_thesis.yaml</code> for fund-specific weighting.</div>',
        )
    thesis_html.append('</section>')

    # --- Section: 5-Dimension ---------------------------------------
    five_d_html: list[str] = ['<section class="card" id="five-d">'
                              '<h2 class="section-title">5-Dimension impact score</h2>']
    if fd:
        rows: list[str] = []
        for dim_name, ds in (
            ("What", fd.what), ("Who", fd.who), ("How Much", fd.how_much),
            ("Contribution", fd.contribution), ("Risk", fd.risk),
        ):
            pct = max(0, min(100, int(ds.score / 5 * 100)))
            rows.append(
                f'<tr><td><b>{dim_name}</b></td>'
                f'<td>{ds.score:.1f}</td>'
                f'<td style="min-width:160px"><div class="bar-track">'
                f'<div class="bar-fill coverage" style="width:{pct}%"></div></div></td>'
                f'<td>{escape(ds.provenance)}</td>'
                f'<td style="font-size:0.88em;color:var(--text-secondary)">{escape(ds.notes or "")}</td>'
                f'</tr>',
            )
        five_d_html.append(
            '<table class="data"><thead><tr>'
            '<th>Dimension</th><th>Score</th><th>Progress</th>'
            '<th>Provenance</th><th>Notes</th></tr></thead><tbody>'
            + "".join(rows) + '</tbody></table>',
        )
        if fd.recommendations:
            items = "".join(f"<li>{escape(r)}</li>" for r in fd.recommendations)
            five_d_html.append(f'<div class="callout"><b>Recommendations</b><ul>{items}</ul></div>')
    else:
        five_d_html.append('<div class="callout warn">No 5D assessment available.</div>')
    five_d_html.append('</section>')

    # --- Section: SDG ----------------------------------------------
    sdg_html: list[str] = ['<section class="card" id="sdg">'
                           '<h2 class="section-title">SDG alignment (top 5)</h2>']
    if sdgs:
        pills = " ".join(sdg_swatch(a.goal, a.score) for a in sdgs[:8])
        sdg_html.append(f'<p class="section-lede">{pills}</p>')
        rows = "".join(
            f'<tr><td>SDG {a.goal} — {escape(a.goal_name)}</td>'
            f'<td>{a.score:.0f}/100</td>'
            f'<td>{len(a.matched_targets)}</td>'
            f'<td style="font-size:0.85em;color:var(--text-secondary)">'
            f'{escape(a.provenance)} ({escape(a.scoring_basis)})</td></tr>'
            for a in sdgs[:5]
        )
        sdg_html.append(
            '<table class="data"><thead><tr><th>SDG</th><th>Score</th>'
            '<th>Targets matched</th><th>Basis</th></tr></thead>'
            f'<tbody>{rows}</tbody></table>',
        )
    else:
        sdg_html.append('<div class="callout">No SDG alignments computed.</div>')
    sdg_html.append('</section>')

    # --- Section: DD + greenwashing -------------------------------
    extras_html: list[str] = ['<section class="card" id="dd-gw">'
                              '<h2 class="section-title">DD & greenwashing</h2>']
    if dd_coverage_pct is not None:
        pct = max(0, min(100, int(dd_coverage_pct)))
        band = "green" if dd_coverage_pct >= 70 else "orange" if dd_coverage_pct >= 40 else "red"
        extras_html.append(
            f'<h3>Due-diligence checklist coverage</h3>'
            f'<div class="bar-track" style="height:14px">'
            f'<div class="bar-fill {band}" style="width:{pct}%"></div></div>'
            f'<p style="margin-top:6px;font-size:0.88em;color:var(--text-secondary)">'
            f'{dd_coverage_pct:.1f}% of checklist items addressed in the document.</p>',
        )
    else:
        extras_html.append('<div class="callout">DD checklist coverage not run.</div>')

    if greenwashing_score is not None:
        kind = gw_kind
        extras_html.append(
            f'<h3>Greenwashing composite</h3>'
            f'<p>Composite risk: <b>{greenwashing_score:.1f}/100</b>'
            + (f' — <span class="pill {kind}">{escape(greenwashing_classification)}</span>'
               if greenwashing_classification else '')
            + '</p>',
        )
    else:
        extras_html.append('<div class="callout">Greenwashing screen not run.</div>')
    extras_html.append('</section>')

    # --- Section: IC Gate detail -----------------------------------
    gate_rows = "".join(
        f'<tr><td>{escape(c.name)}</td>'
        f'<td>{_status_pill(c.status)}</td>'
        f'<td>{escape(str(c.actual))}</td>'
        f'<td>{escape(str(c.threshold))}</td>'
        f'<td style="font-size:0.88em;color:var(--text-secondary)">{escape(c.message)}</td></tr>'
        for c in scorecard.checks
    )
    gate_html: list[str] = [
        '<section class="card" id="gate">'
        '<h2 class="section-title">IC gate detail</h2>',
        '<table class="data"><thead><tr><th>Check</th><th>Status</th>'
        '<th>Actual</th><th>Threshold</th><th>Note</th></tr></thead>'
        f'<tbody>{gate_rows}</tbody></table>',
    ]
    if scorecard.blocking_failures:
        items = "".join(f"<li>{escape(m)}</li>" for m in scorecard.blocking_failures)
        gate_html.append(f'<div class="callout danger"><b>Blocking failures</b><ul>{items}</ul></div>')
    if scorecard.warnings_list:
        items = "".join(f"<li>{escape(m)}</li>" for m in scorecard.warnings_list)
        gate_html.append(f'<div class="callout warn"><b>Warnings</b><ul>{items}</ul></div>')
    gate_html.append('</section>')

    # --- Section: Recommendation ----------------------------------
    rec_html = (
        '<section class="card" id="rec">'
        '<h2 class="section-title">Recommendation</h2>'
        f'<div class="callout ok" style="font-size:1.02em"><b>{scorecard.overall_status.upper()}</b> '
        f'— {escape(scorecard.recommendation)}</div>'
        '</section>'
    )

    toc = render_toc([
        ("thesis-fit", "1. Thesis fit"),
        ("five-d",     "2. 5-Dimension"),
        ("sdg",        "3. SDG alignment"),
        ("dd-gw",      "4. DD & greenwashing"),
        ("gate",       "5. IC gate detail"),
        ("rec",        "6. Recommendation"),
    ])

    body = (
        '<div class="page">'
        f'{toc}'
        '<main>'
        f'{hero}'
        f'{render_kpi_strip(kpis)}'
        f'{"".join(thesis_html)}'
        f'{"".join(five_d_html)}'
        f'{"".join(sdg_html)}'
        f'{"".join(extras_html)}'
        f'{"".join(gate_html)}'
        f'{rec_html}'
        f'{render_footer("Review before submission. Not investment advice.")}'
        '</main>'
        '</div>'
    )
    return wrap_document(
        title=f"IC Memo — {company.name or 'Unnamed'}",
        body_html=body,
    )


def render_ic_memo(
    assessment: Assessment,
    scorecard: DealScorecard,
    thesis: FundThesis | None = None,
    *,
    output_format: Literal["markdown", "html", "docx", "pptx"] = "markdown",
    path: str | Path | None = None,
    **kwargs,
) -> str | Path:
    """Single entry point — returns the rendered string or written file path.

    Supported formats: ``markdown`` (default), ``html``, ``docx``, ``pptx``.
    """
    if output_format == "html":
        html = render_ic_memo_html(assessment, scorecard, thesis, **kwargs)
        if path:
            p = Path(path)
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(html, encoding="utf-8")
            return p
        return html

    md = render_ic_memo_markdown(assessment, scorecard, thesis, **kwargs)
    if output_format == "markdown":
        if path:
            p = Path(path)
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(md, encoding="utf-8")
            return p
        return md
    if output_format == "docx":
        return render_ic_memo_docx(md, path or "ic_memo.docx")
    if output_format == "pptx":
        return render_ic_memo_pptx(assessment, scorecard, thesis, path or "ic_memo.pptx")
    raise ValueError(f"Unknown output_format: {output_format}")
